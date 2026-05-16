"""生成v5汇报PPT（白底黑字，极简风格）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

CHARTS_DIR = '/workspace/xuannv/references/charts'
OUTPUT_PATH = '/workspace/xuannv/references/玄女底座汇报_v5.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 统一颜色 (使用 RGB 元组)
BLACK = (0, 0, 0)
DARK_GRAY = (51, 51, 51)
MED_GRAY = (102, 102, 102)
BLUE = (26, 82, 118)

def set_font(run, size=14, bold=False, color=BLACK, name='Arial'):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = type('obj', (object,), {'__init__': lambda s, r, g, b: None})()
    # 使用正确的方式设置颜色
    from pptx.dml.color import RGBColor
    run.font.color.rgb = RGBColor(*color)
    run.font.name = name

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=BLACK, align=PP_ALIGN.LEFT, font_name='Arial'):
    from pptx.dml.color import RGBColor
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.alignment = align
    return box

def add_image_slide(prs, title, img_file, note=None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    
    add_textbox(slide, 0.5, 0.3, 12, 0.6, title, font_size=28, bold=True, align=PP_ALIGN.LEFT)
    
    img_path = os.path.join(CHARTS_DIR, img_file)
    if os.path.exists(img_path):
        from PIL import Image
        with Image.open(img_path) as img:
            w_px, h_px = img.size
        aspect = h_px / w_px
        img_w = 12.0
        img_h = img_w * aspect
        max_h = 5.0
        if img_h > max_h:
            img_h = max_h
            img_w = img_h / aspect
        img_x = (13.333 - img_w) / 2
        slide.shapes.add_picture(img_path, Inches(img_x), Inches(1.0),
                                  width=Inches(img_w), height=Inches(img_h))
    
    if note:
        add_textbox(slide, 0.5, 6.8, 12, 0.4, note, font_size=11,
                    color=MED_GRAY, align=PP_ALIGN.LEFT)
    
    return slide

# ==================== P1: 封面 ====================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(s1, 0, 2.0, 13.333, 1.2, '玄女底座', font_size=56, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s1, 0, 3.3, 13.333, 0.6, '自研地球嵌入基座 + 可视化展示平台',
            font_size=22, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s1, 0, 4.2, 13.333, 0.5,
            '模型训练 · 下游推理 · 交互标注 · AI 智能体 — 全栈产品',
            font_size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ==================== P2: 模型原理 ====================
add_image_slide(prs, '模型原理：输入 → STP 编码器 → VMF 瓶颈 → 64 维输出',
                '05_model_architecture.png',
                '训练时 Skip L2 防坍缩，推理时 L2 + VMF，时序对比是核心差异。')

# ==================== P3: 核心亮点 ====================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(s3, 0.5, 0.3, 12, 0.6, '核心亮点', font_size=28, bold=True)

highlights = [
    ('01  反坍缩训练', 'Uniformity 为负，嵌入均匀散开。第一次训练崩到 -0.1，调了十几轮才稳住。'),
    ('02  时序敏感', '不重叠双窗口训练，月度级变化可检测。AEF 用年度数据，我们在月度 granularity。'),
    ('03  纯自研', '代码、结构、训练流程全部自己写的。知识产权干净，升级不受制于人。')
]
for i, (title, desc) in enumerate(highlights):
    y = 1.3 + i * 1.8
    add_textbox(s3, 0.8, y, 11.5, 0.5, title, font_size=18, bold=True)
    add_textbox(s3, 0.8, y + 0.5, 11.5, 0.7, desc, font_size=14)

# ==================== P4: 变化检测 ====================
add_image_slide(prs, '变化检测：BA 79.8%（含 CD Head）/ 50%（Raw）',
                '01_cd_comparison_ba.png',
                '加 CD Head 超 AEF 是事实；不加时 50% 也是事实，V6 在补。')

# ==================== P5: 展示平台总览 ====================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(s5, 0.5, 0.3, 12, 0.6, '展示平台总览', font_size=28, bold=True)
add_textbox(s5, 0.5, 1.0, 12, 0.4, '不是 PPT，是真能点、能切、能出报告的网页。',
            font_size=16, bold=True, color=BLUE)

modules = [
    ('三维地球可视化', 'Three.js 全球训练样本分布，缩放到 patch 级别看详情'),
    ('下游监测能力', '5 大 Task Head 动态切换，Mosaic 大图拼接，Patch 详情弹窗'),
    ('AI 智能体报告', '自然语言输入，DeepSeek-V4 解析 → 执行 → 润色出报告')
]
for i, (name, desc) in enumerate(modules):
    y = 1.7 + i * 1.4
    add_textbox(s5, 0.8, y, 3, 0.4, name, font_size=16, bold=True)
    add_textbox(s5, 4.2, y, 8, 0.5, desc, font_size=14)

add_textbox(s5, 0.5, 6.0, 12, 0.4,
            '技术栈：React 19 · TypeScript · Three.js · MapLibre · FastAPI · Docker',
            font_size=11, color=MED_GRAY, align=PP_ALIGN.CENTER)

# ==================== P6: 五大下游任务 ====================
add_image_slide(prs, '五大下游任务：全部已上线',
                '08_downstream_tasks.png',
                '4 / 5 任务纯 CPU 推理，单 patch ~10 ms。CD Head 自动选择空闲 GPU。')

# ==================== P7: 交互式标注训练 ====================
add_image_slide(prs, '交互式标注训练：零代码自定义训练',
                '09_annotation_pipeline.png',
                '业务人员自己就能训专用分类模型；传统方式至少两周，现在 5 分钟。')

# ==================== P8: AI 智能体报告 ====================
add_image_slide(prs, 'AI 智能体报告：DeepSeek-V4 驱动',
                '10_agent_report.png',
                '已经接通，不是规划。输一句话就能出带表格和分析的报告。')

# ==================== P9: 轻量化 ====================
add_image_slide(prs, '轻量化：23 TB → 3 GB → 16 MB',
                '06_compression_ladder.png',
                '16 MB 比表情包还小；平台上所有推理底层都是这 64 字节。')

# ==================== P10: 效率提升 ====================
add_image_slide(prs, '效率提升 vs 传统遥感分析',
                '04_efficiency_comparison.png',
                '数字来自实际项目；存储降是因为下游不需要原始影像。')

# ==================== P11: 升级路线 ====================
add_image_slide(prs, '升级路线：V5 → V6 → V6.5 → V7',
                '07_roadmap.png',
                'V5 已验证，V6 在训，V6.5 已设计，V7 长期目标。')

# ==================== P12: 总结 ====================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(s12, 0.5, 0.3, 12, 0.6, '总结', font_size=28, bold=True)

summary = [
    ('够轻', '64 字节 / km²，23 TB 压到 16 MB，任何设备离线可用'),
    ('够敏', '月度级时序敏感，变化检测 BA 79.8%'),
    ('够全', '模型 + 平台 + 5 大任务 + 交互训练 + AI 智能体，一整套产品'),
    ('够诚', '从零自研，优势不夸大，短板不回避，路线清晰')
]
for i, (word, desc) in enumerate(summary):
    y = 1.3 + i * 1.2
    add_textbox(s12, 1.0, y, 1.5, 0.5, word, font_size=26, bold=True, color=BLUE)
    add_textbox(s12, 3.0, y + 0.05, 9, 0.5, desc, font_size=14)

add_textbox(s12, 0.5, 6.2, 12, 0.4, '有问题随时打断。',
            font_size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

prs.save(OUTPUT_PATH)
print(f'[OK] PPT saved to: {OUTPUT_PATH}')
